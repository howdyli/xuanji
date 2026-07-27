"""PPTX renderer — converts session Markdown to a PowerPoint deck.

Slide mapping (input is the output of ``build_session_markdown``):
- The leading ``# `` heading becomes the title slide (metadata as subtitle).
- Each ``## `` message section becomes one content slide.

Uses ``python-pptx``. If unavailable, raises ``RuntimeError`` with install hint.
"""

from __future__ import annotations

import io
import re

try:
    from pptx import Presentation
    from pptx.util import Pt
except ImportError:  # pragma: no cover - exercised via RuntimeError branch
    Presentation = None  # type: ignore[assignment]

# Keep slides readable: truncate very long message bodies and cap deck size.
_MAX_BODY_CHARS = 1600
_MAX_CONTENT_SLIDES = 60
_BODY_FONT_PT = 14

_MD_NOISE_RE = re.compile(r"(\*\*|__|`{1,3})")
_FOOTER_RE = re.compile(r"^\*.+\*$")


def _strip_md(text: str) -> str:
    """Light-weight Markdown → plain text for slide bodies."""
    return _MD_NOISE_RE.sub("", text).strip()


def _split_sections(
    markdown_text: str,
) -> tuple[str, str, list[tuple[str, str]]]:
    """Parse session Markdown into (title, metadata, [(heading, body), ...])."""
    title: str | None = None
    meta = ""
    sections: list[tuple[str, str]] = []
    heading: str | None = None
    body: list[str] = []

    for line in markdown_text.splitlines():
        stripped = line.strip()
        if stripped == "---" or _FOOTER_RE.match(stripped):
            continue
        if stripped.startswith("## "):
            if heading is not None:
                sections.append((heading, "\n".join(body).strip()))
            heading = stripped[3:].strip()
            body = []
        elif heading is not None:
            body.append(line)
        elif stripped.startswith("# ") and title is None:
            title = stripped[2:].strip()
        elif stripped.startswith(">") and not meta:
            meta = stripped.lstrip("> ").strip()

    if heading is not None:
        sections.append((heading, "\n".join(body).strip()))

    return title or "会话导出", meta, sections


def render_markdown_to_pptx(markdown_text: str) -> bytes:
    """Render *markdown_text* to a PPTX byte stream.

    Raises
    ------
    RuntimeError
        If python-pptx is not available.
    """
    if Presentation is None:
        raise RuntimeError(
            "python-pptx not available. Install with:\n    pip install python-pptx"
        )

    title, meta, sections = _split_sections(markdown_text)

    prs = Presentation()

    # ── Title slide ──
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if meta:
        try:
            slide.placeholders[1].text = _strip_md(meta)
        except KeyError:  # pragma: no cover - default template has a subtitle
            pass

    # ── One content slide per message section ──
    truncated = len(sections) > _MAX_CONTENT_SLIDES
    for heading, body in sections[:_MAX_CONTENT_SLIDES]:
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = _strip_md(heading)

        text = _strip_md(body)
        if len(text) > _MAX_BODY_CHARS:
            text = text[:_MAX_BODY_CHARS] + "…（内容过长已截断）"

        tf = s.placeholders[1].text_frame
        tf.word_wrap = True
        lines = text.split("\n") or [""]
        tf.text = lines[0]
        for extra in lines[1:]:
            tf.add_paragraph().text = extra
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.size = Pt(_BODY_FONT_PT)

    if truncated:
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = "内容过长已截断"
        s.placeholders[1].text_frame.text = (
            f"本会话共 {len(sections)} 条消息，"
            f"仅导出前 {_MAX_CONTENT_SLIDES} 条。完整内容请使用 PDF/Markdown 导出。"
        )

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
