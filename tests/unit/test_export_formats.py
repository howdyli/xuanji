"""Tests for the new export formats: HTML renderer, PPTX renderer, ExportService integration."""

from __future__ import annotations

import io

import pytest
from unittest.mock import MagicMock, AsyncMock
from dataclasses import dataclass

from xiaopaw.export.html_renderer import render_markdown_to_html
from xiaopaw.export.markdown_builder import build_session_markdown
from xiaopaw.export.pptx_renderer import (
    render_markdown_to_pptx,
    _split_sections,
    _strip_md,
    _MAX_CONTENT_SLIDES,
)


# ── Fixtures ────────────────────────────────────────────────────────────────

SAMPLE_MESSAGES = [
    {"role": "user", "content": "你好，帮我写一份报告", "ts": 1720000000000, "feishu_msg_id": None},
    {"role": "assistant", "content": "好的，这是一份**报告**。", "ts": 1720000060000, "feishu_msg_id": None},
]


def _session_md() -> str:
    return build_session_markdown("测试会话", "s-test-001", "2026-07-08", SAMPLE_MESSAGES)


# ═══════════════════════════════════════════════════════════════════════════
# HTML Renderer Tests
# ═══════════════════════════════════════════════════════════════════════════


class TestHtmlRenderer:
    def test_generates_standalone_html(self):
        """输出为完整独立 HTML 文档（含 DOCTYPE + zh-CN + style）"""
        result = render_markdown_to_html(_session_md())
        assert isinstance(result, bytes)
        text = result.decode("utf-8")
        assert text.startswith("<!DOCTYPE html>")
        assert 'lang="zh-CN"' in text
        assert "<style>" in text

    def test_content_rendered(self):
        """会话标题与消息内容出现在 HTML 中"""
        text = render_markdown_to_html(_session_md()).decode("utf-8")
        assert "测试会话" in text
        assert "你好，帮我写一份报告" in text

    def test_markdown_converted_to_html_tags(self):
        """Markdown 语法被转换为 HTML 标签（H1 / 加粗 / 表格）"""
        md = "# Title\n\n**bold**\n\n| A | B |\n|---|---|\n| 1 | 2 |"
        text = render_markdown_to_html(md).decode("utf-8")
        assert "<h1>Title</h1>" in text
        assert "<strong>bold</strong>" in text
        assert "<table>" in text

    def test_chinese_content_utf8(self):
        """中文内容以 UTF-8 编码无损输出"""
        md = "# 中文标题\n\n这是一段中文内容。"
        result = render_markdown_to_html(md)
        assert "这是一段中文内容。" in result.decode("utf-8")


# ═══════════════════════════════════════════════════════════════════════════
# PPTX Renderer Tests
# ═══════════════════════════════════════════════════════════════════════════


class TestPptxSectionParsing:
    def test_split_sections_basic(self):
        """标题 / 元数据 / 消息分节解析正确"""
        title, meta, sections = _split_sections(_session_md())
        assert title == "测试会话"
        assert "s-test-001" in meta
        assert len(sections) == 2
        assert "User" in sections[0][0]
        assert "你好，帮我写一份报告" in sections[0][1]

    def test_split_sections_skips_hr_and_footer(self):
        """--- 分隔线与 *页脚* 不进入正文"""
        _, _, sections = _split_sections(_session_md())
        for _, body in sections:
            assert "---" not in body
            assert "玄机" not in body

    def test_split_sections_empty_input(self):
        """空输入返回默认标题与空分节"""
        title, meta, sections = _split_sections("")
        assert title == "会话导出"
        assert meta == ""
        assert sections == []

    def test_strip_md_removes_noise(self):
        """加粗/行内代码标记被清除"""
        assert _strip_md("**bold** and `code` and __em__") == "bold and code and em"


class TestPptxRenderer:
    def test_generates_valid_pptx(self):
        """生成的字节以 PK 开头（ZIP 格式）"""
        result = render_markdown_to_pptx(_session_md())
        assert isinstance(result, bytes)
        assert result[:2] == b"PK"
        assert len(result) > 1000

    def test_slide_count_matches_messages(self):
        """幻灯片数 = 1 标题页 + 每条消息一页"""
        from pptx import Presentation

        result = render_markdown_to_pptx(_session_md())
        prs = Presentation(io.BytesIO(result))
        assert len(prs.slides) == 1 + len(SAMPLE_MESSAGES)

    def test_title_and_content_in_slides(self):
        """标题页与内容页文本正确"""
        from pptx import Presentation

        prs = Presentation(io.BytesIO(render_markdown_to_pptx(_session_md())))
        all_text = "\n".join(
            shape.text_frame.text
            for slide in prs.slides
            for shape in slide.shapes
            if shape.has_text_frame
        )
        assert "测试会话" in all_text
        assert "你好，帮我写一份报告" in all_text
        # Markdown 加粗标记已被清除
        assert "**" not in all_text

    def test_deck_size_capped(self):
        """超过上限的消息被截断并追加提示页"""
        from pptx import Presentation

        many = [
            {"role": "user", "content": f"消息 {i}", "ts": 1720000000000 + i, "feishu_msg_id": None}
            for i in range(_MAX_CONTENT_SLIDES + 10)
        ]
        md = build_session_markdown("长会话", "s-long", "2026-07-08", many)
        prs = Presentation(io.BytesIO(render_markdown_to_pptx(md)))
        # 1 标题页 + 上限内容页 + 1 截断提示页
        assert len(prs.slides) == 1 + _MAX_CONTENT_SLIDES + 1

    def test_raises_without_pptx(self):
        """python-pptx 不可用时抛出 RuntimeError"""
        from xiaopaw.export import pptx_renderer

        original = pptx_renderer.Presentation
        try:
            pptx_renderer.Presentation = None
            with pytest.raises(RuntimeError, match="python-pptx not available"):
                render_markdown_to_pptx("# test")
        finally:
            pptx_renderer.Presentation = original


# ═══════════════════════════════════════════════════════════════════════════
# ExportService — pptx / html 集成
# ═══════════════════════════════════════════════════════════════════════════


@dataclass(frozen=True)
class FakeSessionEntry:
    id: str = "s-test"
    title: str = "测试会话"
    created_at: str = "2026-07-08"
    verbose: bool = False
    message_count: int = 2


@dataclass(frozen=True)
class FakeMessageEntry:
    role: str = "user"
    content: str = "你好"
    ts: int = 1720000000000
    feishu_msg_id: str | None = None


@pytest.fixture
def mock_session_mgr():
    mgr = MagicMock()
    mgr.get_session_by_id = AsyncMock(return_value=FakeSessionEntry())
    mgr.load_history = AsyncMock(
        return_value=[
            FakeMessageEntry(role="user", content="你好", ts=1720000000000),
            FakeMessageEntry(role="assistant", content="你好！有什么可以帮你？", ts=1720000060000),
        ]
    )
    return mgr


class TestExportServiceNewFormats:
    @pytest.mark.asyncio
    async def test_export_pptx(self, mock_session_mgr):
        """PPTX 导出返回正确的文件名与 content_type"""
        from xiaopaw.export.service import ExportService

        svc = ExportService(mock_session_mgr)
        file_bytes, filename, content_type = await svc.export_session("s-test", "pptx")

        assert filename == "s-test.pptx"
        assert "presentationml" in content_type
        assert file_bytes[:2] == b"PK"

    @pytest.mark.asyncio
    async def test_export_html(self, mock_session_mgr):
        """HTML 导出返回正确的文件名与 content_type"""
        from xiaopaw.export.service import ExportService

        svc = ExportService(mock_session_mgr)
        file_bytes, filename, content_type = await svc.export_session("s-test", "html")

        assert filename == "s-test.html"
        assert content_type == "text/html; charset=utf-8"
        text = file_bytes.decode("utf-8")
        assert text.startswith("<!DOCTYPE html>")
        assert "测试会话" in text

    @pytest.mark.asyncio
    async def test_unsupported_format_message_lists_new_formats(self, mock_session_mgr):
        """错误消息包含全部支持的格式"""
        from xiaopaw.export.service import ExportService

        svc = ExportService(mock_session_mgr)
        with pytest.raises(ValueError, match="pptx"):
            await svc.export_session("s-test", "xlsx")
